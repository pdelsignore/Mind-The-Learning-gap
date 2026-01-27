#!/usr/bin/env python3
"""
Script to create Mind the Learning Gap presentation
Run: pip install python-pptx && python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN

# Pearson colors
PEARSON_NAVY = RgbColor(0x0D, 0x00, 0x4D)
PEARSON_PURPLE = RgbColor(0x51, 0x2E, 0xAB)
PEARSON_YELLOW = RgbColor(0xFF, 0xCE, 0x00)

# Slide content
slides_content = [
    {
        "title": "Mind the Learning Gap",
        "subtitle": "The Missing Link in AI's Productivity Promise"
    },
    {
        "title": "The Size of the Prize",
        "content": "$4.8T – $6.6T\nPotential US economic impact by 2034\n\nUnlocking this prize demands learning and augmentation, together."
    },
    {
        "title": "The Human Imperative",
        "content": '"Every positive scenario for an AI-enabled future is built on human development. Every negative one stems from neglecting it."\n\n— Omar Abbosh, CEO, Pearson'
    },
    {
        "title": "The D.E.E.P. Framework",
        "content": "A roadmap to unlock AI productivity through learning:\n\n• Diagnose — Define your task augmentation plan\n• Embed — Instill learning in the flow of work\n• Evaluate — Measure progress with skills assessment\n• Prioritize — Position learning as strategic investment"
    },
    {
        "title": "Diagnose",
        "content": "Define your task augmentation plan\n\n• Conduct task-based analysis\n• Identify your expert enthusiasts\n• Build your augmentation squads\n• Identify and roll out use cases"
    },
    {
        "title": "Embed",
        "content": "Instill learning in the flow of work\n\n• Create a culture of continuous learning\n• Embed learning in daily workflows\n• Enable social learning\n• Emphasize durable skills"
    },
    {
        "title": "Evaluate",
        "content": "Measure progress with robust assessment\n\n• Build usable skills data infrastructure\n• Invest in ambient assessment methods\n• Use AI to personalize learning\n• Test skills in authentic conditions"
    },
    {
        "title": "Prioritize",
        "content": "Position learning as strategic investment\n\n• Redefine L&D as capability curators\n• Prioritize skills over roles\n• Build a measurable skills ecosystem\n• Incentivize continuous learning"
    },
    {
        "title": "The Augmented Knowledge Worker",
        "content": "Six ways AI will enhance human work:\n\n• Knowledge Codification\n• Data Gathering & Ideation\n• Multi-Agent Collaboration\n• Rapid & Continuous Learning\n• Personalized Coaching\n• Seamless Human-AI Teamwork"
    },
    {
        "title": "Economic Impact Scenarios",
        "content": "US Economy by 2034:\n\n$4.8T — Conservative (7% elasticity)\n$5.4T — Moderate (10% elasticity)\n$6.6T — Optimistic (high adoption)"
    },
    {
        "title": "Close the Learning Gap",
        "content": "The future depends on how we invest in people and learning.\n\nHuman development is the foundation of progress."
    }
]

def create_presentation():
    # Try to use template, fall back to blank
    try:
        prs = Presentation('Pearson_PPT-Template_MASTER.potx')
        # Remove template slides
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    except:
        prs = Presentation()

    # Get slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_layout = prs.slide_layouts[1]  # Title and content

    for i, slide_data in enumerate(slides_content):
        if i == 0:
            # Title slide
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = slide_data["title"]
            subtitle.text = slide_data["subtitle"]
        else:
            # Content slides
            slide = prs.slides.add_slide(content_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = slide_data["title"]
            body.text = slide_data["content"]

    # Save
    output_path = 'Mind-the-Learning-Gap-Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
